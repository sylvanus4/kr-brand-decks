#!/usr/bin/env python3
"""fill_template.py -- fill an EXISTING .pptx template with your content, preserving
its design 100%.

Inspired by ppt-master's highest-ROI idea: don't rebuild the pptx from scratch --
clone the original and patch only the text. python-pptx edits text runs in place, so
every color, font, layout, image, and animation in your template is kept exactly; only
the {{placeholder}} tokens (or literal strings you map) are replaced.

Use when you already have a branded company .pptx and just want to drop in new numbers/
text -- the from-scratch engine (_engine/render_deck.py) is for building a deck when you
have NO template.

Two modes:
  1) Token mode  -- template has {{key}} placeholders; values.json maps key -> text.
  2) Map mode    -- values.json maps an exact existing string -> replacement string.

Usage:
  python3 tools/fill_template.py --in template.pptx --values values.json --out filled.pptx
  python3 tools/fill_template.py --in template.pptx --list-text     # dump editable text to author your map
"""
import argparse
import json
import re

from pptx import Presentation

TOKEN = re.compile(r"\{\{\s*([^}]+?)\s*\}\}")


def iter_runs(prs):
    for si, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        yield si, r
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                yield si, r


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="inp", required=True)
    ap.add_argument("--values", help="JSON: {key_or_string: replacement}")
    ap.add_argument("--out")
    ap.add_argument("--list-text", action="store_true", help="print editable text + tokens, then exit")
    args = ap.parse_args()

    prs = Presentation(args.inp)

    if args.list_text:
        seen_tokens = set()
        for si, r in iter_runs(prs):
            if r.text.strip():
                print(f"[slide {si + 1}] {r.text}")
            seen_tokens.update(TOKEN.findall(r.text))
        if seen_tokens:
            print("\nTokens found:", ", ".join(sorted(seen_tokens)))
        return

    if not args.values or not args.out:
        ap.error("--values and --out are required unless --list-text")
    with open(args.values, encoding="utf-8") as f:
        values = json.load(f)

    n_tok, n_map = 0, 0
    for si, r in iter_runs(prs):
        txt = r.text
        # token mode: {{key}} -> values[key]
        def _sub(m):
            nonlocal n_tok
            k = m.group(1).strip()
            if k in values:
                n_tok += 1
                return str(values[k])
            return m.group(0)
        new = TOKEN.sub(_sub, txt)
        # map mode: exact-string replacement (only if no token changed this run)
        if new == txt:
            for k, v in values.items():
                if k in new and not k.startswith("{{"):
                    new = new.replace(k, str(v)); n_map += 1
        if new != txt:
            r.text = new

    prs.save(args.out)
    print(f"[ok] {args.out}  (token replacements: {n_tok}, string replacements: {n_map})")
    print("Design preserved: only text runs were patched; all formatting/layout/images kept.")


if __name__ == "__main__":
    main()
