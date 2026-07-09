#!/usr/bin/env python3
"""validate.py -- deterministic quality gate for a rendered deck.

Exit 0 = PASS, non-zero = FAIL. Checks: file opens, expected slide count,
the brand accent color actually appears in the deck XML, and no leftover
placeholder/lorem text. Run before publishing any example deck.

Usage:
  python3 validate.py deck.pptx --palette palette.json --expect-slides 6
"""
import argparse
import json
import sys
import re
from pptx import Presentation

LOREM = re.compile(r"lorem ipsum|\{\{.*?\}\}|\bTODO\b|xxxx", re.IGNORECASE)


def _lum(h):
    h = h.lstrip("#")
    r, g, b = (int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))
    f = lambda c: c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * f(r) + 0.7152 * f(g) + 0.0722 * f(b)


def _contrast(a, b):
    la, lb = _lum(a), _lum(b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def check_contrast(pal):
    """Palette readability gate: cover text on the divider bg, eyebrow/accent text on
    the body bg, and the accent tick vs the divider bg (would catch an invisible tick)."""
    fails = []
    dbg, dink = pal["divider_bg"], pal.get("divider_ink", "FFFFFF")
    acc = pal["accent"]
    acc_ink = pal.get("accent_ink", acc)
    bg = pal.get("bg", "FFFFFF")
    if _contrast(dink, dbg) < 3.0:
        fails.append(f"cover/divider text {dink} on {dbg} = {_contrast(dink, dbg):.1f}:1 (<3.0)")
    if _contrast(acc_ink, bg) < 2.4:
        fails.append(f"eyebrow/accent text {acc_ink} on {bg} = {_contrast(acc_ink, bg):.1f}:1 (<2.4)")
    tick = dink if dbg.lstrip("#").upper() == acc.lstrip("#").upper() else acc
    if _contrast(tick, dbg) < 1.25:
        fails.append(f"accent tick {tick} nearly invisible on {dbg} ({_contrast(tick, dbg):.2f}:1)")
    return fails


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--palette")
    ap.add_argument("--expect-slides", type=int)
    ap.add_argument("--check-contrast", action="store_true",
                    help="also gate palette readability (cover text, eyebrow, accent tick)")
    args = ap.parse_args()

    fails = []
    try:
        prs = Presentation(args.deck)
    except Exception as e:
        print(f"FAIL: cannot open deck: {e}")
        sys.exit(1)

    n = len(prs.slides._sldIdLst)
    if args.expect_slides and n != args.expect_slides:
        fails.append(f"slide count {n} != expected {args.expect_slides}")

    # gather all text
    texts = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    blob = "\n".join(texts)
    if LOREM.search(blob):
        fails.append("placeholder/lorem text found")
    if not blob.strip():
        fails.append("deck has no text")

    # accent color present in XML?
    if args.palette:
        with open(args.palette, encoding="utf-8") as f:
            pal = json.load(f)
        accent = pal.get("accent", "").lstrip("#").upper()
        import zipfile
        xml = ""
        with zipfile.ZipFile(args.deck) as z:
            for nm in z.namelist():
                if nm.startswith("ppt/slides/slide") and nm.endswith(".xml"):
                    xml += z.read(nm).decode("utf-8", "ignore").upper()
        if accent and accent not in xml:
            fails.append(f"accent {accent} not found in slide XML")
        if args.check_contrast:
            fails.extend(check_contrast(pal))

    if fails:
        print("VERDICT: FAIL")
        for f in fails:
            print("  - " + f)
        sys.exit(1)
    print(f"VERDICT: PASS ({n} slides, accent applied, no lorem)")


if __name__ == "__main__":
    main()
