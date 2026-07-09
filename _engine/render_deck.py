#!/usr/bin/env python3
"""render_deck.py -- palette-driven native PPTX generator (no template).

Reads a brand PALETTE json (company colors + font) and a content SPEC json, then
builds a deck of reference-grade slides natively with python-pptx. The model writes
CONTENT; this code owns all FORMAT (color, font, layout) -- every color is pulled
from the palette so each company's deck is automatically on-brand.

Standalone by design: no proprietary template, no hardcoded interpreter, no repo
paths, no external services. Optional chart slides use charts.py (degrade to text
if matplotlib/KR font missing). Optional hero images are NOT fetched here.

Usage:
  python3 render_deck.py --palette PALETTE.json --spec SPEC.json --out deck.pptx
  python3 render_deck.py ... --pdf        # also export PDF via LibreOffice (soffice)

Palette roles: ink sub muted surface border bg accent divider_bg divider_ink (+ font).
Spec: {"meta": {...theme...}, "slides": [ {"layout": "...", ...}, ... ]}
Layouts: cover · toc · divider · icongrid · textfigure · table · numbered · bullets ·
  roadmap · kpi · bignum · trend · segment · diagram (flow/nodes architecture) ·
  imagefull · imagesplit · imagegrid · closing. Visual theme via --theme / meta.theme.
"""
import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_IN = 914400
SW, SH = 13.333, 7.5


def C(h):
    return RGBColor.from_string(h.lstrip("#").upper())


# A THEME is the VISUAL STYLE (typography, cover/divider treatment, accent, density,
# emoji) -- orthogonal to the brand PALETTE (colors). Defaults reproduce the original
# "corporate" look, so a deck with no theme renders exactly as before.
THEME_DEFAULTS = {
    "cover_style": "solid",       # solid | minimal | band | sidebar | dark
    "cover_title_pt": 54,
    "cover_align": "left",        # left | center
    "title_pt": 25,               # content-slide header title
    "eyebrow_track": 120,         # letter-spacing on the eyebrow
    "eyebrow_upper": True,        # UPPERCASE the eyebrow (Latin only)
    "divider_style": "solid",     # solid | minimal | huge
    "divider_num_pt": 60,
    "accent_style": "tick",       # tick | underline | block | none
    "content_bg": "bg",           # bg | surface  (content-slide background token)
    "hairline": True,             # header hairline under the title
    "subtitle_pt": 14,
    "body_pt": 12,
    "emoji": False,               # render a leading emoji on eyebrows/sections
    "section_icons": False,       # a Lucide section icon beside content-slide titles
    "card_icons": True,           # a Lucide icon per bignum KPI card
    "closing_mark": True,         # a Lucide mark on the closing slide
}

# Korean section keyword -> Lucide icon (line icons, NOT emoji). Used to place real
# icons next to titles / on KPI cards. First substring match wins.
SECTION_ICON = {
    "개요": "layout-grid", "사업": "briefcase", "전략": "compass", "성장": "trending-up",
    "로드맵": "map", "기술": "cpu", "리더십": "users", "제품": "package", "시장": "globe",
    "미래": "rocket", "비전": "telescope", "실적": "chart-bar", "재무": "landmark",
    "핵심": "key", "지표": "gauge", "포트폴리오": "layers", "경쟁": "swords",
    "고객": "user-check", "조직": "building-2", "운영": "settings", "품질": "badge-check",
    "보안": "shield", "데이터": "database", "AI": "bot", "글로벌": "earth",
    "투자": "piggy-bank", "리스크": "triangle-alert", "성과": "trophy", "일정": "calendar",
    "요약": "file-text", "결론": "flag", "배경": "book-open", "문제": "circle-alert",
    "해결": "lightbulb", "방법": "wrench", "사례": "folder-open", "파트너": "handshake",
    "채널": "send", "브랜드": "sparkles", "매출": "trending-up", "이익": "chart-bar",
    "현금": "landmark", "부문": "layers", "건전성": "shield-check",
}


def resolve_icon(text):
    if not text:
        return None
    for kw, ic in SECTION_ICON.items():
        if kw in text:
            return ic
    return None

# eyebrow/section keyword -> emoji (only used when the theme opts in). Formal decks
# (finance/IR/gov/data-report) keep emoji off.
EMOJI_MAP = {
    "개요": "🚀", "사업": "🏢", "전략": "🎯", "성장": "📈", "로드맵": "🗺️",
    "기술": "⚙️", "리더십": "🏆", "제품": "📦", "시장": "🌐", "미래": "✨",
    "비전": "🌟", "실적": "📊", "재무": "💰", "핵심": "🔑", "지표": "📌",
    "portfolio": "🧩", "growth": "📈", "contents": "📑", "business": "🏢",
}


class Deck:
    def __init__(self, pal, tmpdir, theme=None):
        self.p = pal
        self.t = {**THEME_DEFAULTS, **(theme or {})}
        self.font = pal.get("font", "Pretendard")
        self.tmp = tmpdir
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(SW * EMU_IN))
        self.prs.slide_height = Emu(int(SH * EMU_IN))
        self.blank = self.prs.slide_layouts[6]
        self._n = 0
        # persistent icon cache: fetch each Lucide icon once, reuse across runs
        # (faster iteration + offline rebuilds after the first fetch).
        self.icons_dir = os.path.expanduser("~/.cache/kr-brand-decks/icons")
        # accent_ink = accent used as TEXT on a light bg (dark-enough variant for
        # bright brands like yellow/green). on_accent = text drawn ON an accent fill.
        self.accent_ink = pal.get("accent_ink", pal["accent"])
        self.on_accent = pal.get("on_accent", pal["divider_ink"])

    # ---- primitives ----
    def slide(self, bg):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C(bg)
        return s

    def box(self, s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        return tf

    def para(self, tf, text, size, color, bold=False, first=False, sb=0, sa=0,
             align=PP_ALIGN.LEFT, ls=1.15, tracking=None):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(sb); p.space_after = Pt(sa); p.line_spacing = ls
        r = p.add_run(); r.text = text
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = C(color); f.name = self.font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea)
        ea.set("typeface", self.font)
        if tracking is not None:
            rPr.set("spc", str(int(tracking)))
        return p

    def rect(self, s, x, y, w, h, color, line=None):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb = C(color)
        if line:
            r.line.color.rgb = C(line); r.line.width = Pt(1)
        else:
            r.line.fill.background()
        r.shadow.inherit = False
        return r

    def _bar_on(self, bg):
        """Accent tick color that stays visible: white when bg equals the accent."""
        a = self.p["accent"].lstrip("#").upper()
        return self.p["divider_ink"] if bg.lstrip("#").upper() == a else self.p["accent"]

    def pic_contain(self, s, png, x, y, w, h):
        from PIL import Image
        iw, ih = Image.open(png).size
        bx, ix = w / h, iw / ih
        if ix > bx:
            nw, nh = w, w / ix
        else:
            nh, nw = h, h * ix
        s.shapes.add_picture(png, Inches(x + (w - nw) / 2), Inches(y + (h - nh) / 2),
                             Inches(nw), Inches(nh))

    def pic_cover(self, s, png, x, y, w, h):
        """Fill the box with the image, cropping the overflow (cover-fit)."""
        from PIL import Image
        iw, ih = Image.open(png).size
        box_ar, img_ar = w / h, iw / ih
        pic = s.shapes.add_picture(png, Inches(x), Inches(y), Inches(w), Inches(h))
        if img_ar > box_ar:
            crop = (1 - box_ar / img_ar) / 2
            pic.crop_left = crop; pic.crop_right = crop
        else:
            crop = (1 - img_ar / box_ar) / 2
            pic.crop_top = crop; pic.crop_bottom = crop
        return pic

    def _img_or_placeholder(self, s, sp, x, y, w, h):
        """cover-fit sp['image'] into the box, or draw a labelled placeholder panel."""
        img = sp.get("image")
        if img and os.path.exists(img):
            self.pic_cover(s, img, x, y, w, h)
            return True
        self.rect(s, x, y, w, h, self.p["surface"])
        tf = self.box(s, x + 0.3, y + h / 2 - 0.3, w - 0.6, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        self.para(tf, "[image]\n" + sp.get("image_prompt", sp.get("title", "")), 11,
                  self.p["muted"], first=True, align=PP_ALIGN.CENTER)
        return False

    # ---- theme helpers ----
    def cslide(self):
        """Content-slide background per theme + theme decoration."""
        s = self.slide(self.p.get(self.t["content_bg"], self.p["bg"]))
        self.theme_decorate(s, "content")
        return s

    def dslide(self, bg, where):
        """Create a slide with bg and apply the theme's signature decoration."""
        s = self.slide(bg)
        self.theme_decorate(s, where)
        return s

    def theme_decorate(self, s, where):
        """Signature decorative details per theme (behind content). `where` in
        {content, cover, divider}. No-op unless the theme sets `decor`."""
        d = self.t.get("decor")
        if not d:
            return
        pal = self.p
        dark = where in ("cover", "divider")
        if d == "blueprint":
            gc = pal["surface"] if dark else pal["border"]
            for gx in range(1, 13):
                self.rect(s, gx * (SW / 13), 0, 0.006, SH, gc)
            for gy in range(1, 8):
                self.rect(s, 0, gy * (SH / 8), SW, 0.006, gc)
            for cx in (0.3, SW - 0.42):
                for cy in (0.3, SH - 0.42):
                    self.rect(s, cx, cy + 0.055, 0.12, 0.012, pal["accent"])
                    self.rect(s, cx + 0.055, cy, 0.012, 0.12, pal["accent"])
        elif d == "swiss":
            col = pal["divider_ink"] if dark else pal["border"]
            for gx in (SW * 0.33, SW * 0.5, SW * 0.67):
                self.rect(s, gx, 0.3, 0.006, SH - 0.6, col)
        elif d == "neon" and dark:
            for (x, y, w, h) in [(0.35, 0.35, SW - 0.7, 0.02), (0.35, SH - 0.37, SW - 0.7, 0.02),
                                 (0.35, 0.35, 0.02, SH - 0.7), (SW - 0.37, 0.35, 0.02, SH - 0.7)]:
                self.rect(s, x, y, w, h, pal["accent"])
        elif d == "brutalist" and dark:
            # thick accent bars (visible on both the white band-top and the color band)
            self.rect(s, 0.0, 0.0, SW, 0.14, pal["accent"])
            self.rect(s, 0.0, SH - 0.14, SW, 0.14, pal["accent"])
        elif d == "luxury":
            fc = pal["accent"] if dark else pal["border"]
            for (x, y, w, h) in [(0.5, 0.5, SW - 1.0, 0.012), (0.5, SH - 0.51, SW - 1.0, 0.012),
                                 (0.5, 0.5, 0.012, SH - 1.0), (SW - 0.51, 0.5, 0.012, SH - 1.0)]:
                self.rect(s, x, y, w, h, fc)
        elif d == "terminal" and dark:
            for i, cc in enumerate((pal["accent"], pal["muted"], pal["divider_ink"])):
                dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5 + i * 0.28), Inches(0.42),
                                         Inches(0.16), Inches(0.16))
                dot.fill.solid(); dot.fill.fore_color.rgb = C(cc); dot.line.fill.background()
                dot.shadow.inherit = False
            tf = self.box(s, 0.5, 6.72, 11, 0.4)
            self.para(tf, "user@deck:~$ present ▮", 13, pal["accent"], first=True, tracking=20)
        elif d == "newsprint" and not dark:
            self.rect(s, 0.5, 1.7, SW - 1.0, 0.02, pal["ink"])
            self.rect(s, SW * 0.5, 2.3, 0.006, SH - 3.0, pal["border"])

    def _eyebrow_text(self, eyebrow):
        t = self.t
        if t["emoji"] and eyebrow:
            key = next((k for k in EMOJI_MAP if k in eyebrow), None)
            em = EMOJI_MAP.get(key or "", "")
            if em:
                eyebrow = f"{em} {eyebrow}"
        return eyebrow.upper() if t["eyebrow_upper"] else eyebrow

    def _header_accent(self, s):
        pal, a = self.p, self.t["accent_style"]
        y = 2.4
        if a == "none":
            if self.t["hairline"]:
                self.rect(s, 0.62, y + 0.02, 12.1, 0.012, pal["border"])
            return
        if a == "underline":
            self.rect(s, 0.62, y + 0.01, 12.1, 0.03, pal["accent"])
            return
        if a == "block":
            self.rect(s, 0.62, y - 0.02, 0.52, 0.1, pal["accent"])
            if self.t["hairline"]:
                self.rect(s, 1.26, y + 0.03, 11.46, 0.012, pal["border"])
            return
        self.rect(s, 0.62, y, 1.1, 0.045, pal["accent"])              # tick (default)
        if self.t["hairline"]:
            self.rect(s, 1.72, y + 0.018, 11.0, 0.012, pal["border"])

    # ---- shared header for content slides ----
    def header(self, s, eyebrow, title, subtitle, page):
        pal, t = self.p, self.t
        if eyebrow:
            tf = self.box(s, 0.62, 0.44, 9.5, 0.32)
            self.para(tf, self._eyebrow_text(eyebrow), 11, self.accent_ink, bold=True,
                      first=True, tracking=t["eyebrow_track"])
        if page:
            tf = self.box(s, SW - 1.5, 0.44, 0.9, 0.32)
            self.para(tf, page, 11, pal["muted"], first=True, align=PP_ALIGN.RIGHT)
        tx = 0.62
        if t["section_icons"]:
            ic = resolve_icon(eyebrow) or resolve_icon(title)
            if ic and self._icon(s, ic, 0.62, 0.96, 0.42, pal["muted"]):
                tx = 1.2
        tf = self.box(s, tx, 0.92, 12.72 - tx, 1.0)
        self.para(tf, title, t["title_pt"], pal["ink"], bold=True, first=True)
        if subtitle:
            tf = self.box(s, 0.62, 1.78, 12.1, 0.6)
            self.para(tf, subtitle, t["subtitle_pt"], pal["sub"], first=True)
        self._header_accent(s)

    # ---- layouts ----
    def cover(self, sp):
        pal, t = self.p, self.t
        style = t["cover_style"]
        ey = sp.get("eyebrow", "")
        ey = ey.upper() if t["eyebrow_upper"] else ey
        title, sub = sp["title"], sp.get("subtitle")
        ctpt = t["cover_title_pt"]
        al = PP_ALIGN.CENTER if t["cover_align"] == "center" else PP_ALIGN.LEFT
        cx = 0.62

        # AI-generated template background image (creative-template pipeline)
        if sp.get("bg_image") and os.path.exists(sp["bg_image"]):
            self.imagefull({"image": sp["bg_image"], "eyebrow": sp.get("eyebrow", ""),
                            "title": title, "subtitle": sub})
            return

        if style == "minimal":  # light bg, ink text, small accent — editorial/luxury
            s = self.dslide(pal["bg"], "cover")
            self.rect(s, cx, 0.72, 0.9, 0.09, pal["accent"])
            tf = self.box(s, cx, 1.0, 11.5, 0.5)
            self.para(tf, ey, 13, self.accent_ink, bold=True, first=True, tracking=180, align=al)
            tf = self.box(s, cx, 3.7, 12.1, 2.8)
            self.para(tf, title, ctpt, pal["ink"], bold=True, first=True, ls=1.05, align=al)
            if sub:
                self.para(tf, sub, 16, pal["sub"], sb=14, ls=1.25, align=al)
            return

        if style == "band":  # light top + brand-color bottom band holding the title
            s = self.dslide(pal["bg"], "cover")
            bh = 3.0
            self.rect(s, 0, SH - bh, SW, bh, pal["divider_bg"])
            tf = self.box(s, cx, 0.95, 11.5, 0.5)
            self.para(tf, ey, 13, self.accent_ink, bold=True, first=True, tracking=180, align=al)
            tf = self.box(s, cx, SH - bh + 0.55, 12.1, bh - 0.7)
            self.para(tf, title, min(ctpt, 46), pal["divider_ink"], bold=True, first=True, ls=1.05, align=al)
            if sub:
                self.para(tf, sub, 15, pal["divider_ink"], sb=12, ls=1.25, align=al)
            return

        if style == "sidebar":  # left color strip + light main area, ink title
            s = self.dslide(pal["bg"], "cover")
            self.rect(s, 0, 0, 3.3, SH, pal["divider_bg"])
            self.rect(s, 0.6, 0.7, 0.9, 0.12, self._bar_on(pal["divider_bg"]))
            tf = self.box(s, 0.6, 1.0, 2.4, 3.0)
            self.para(tf, ey, 12, pal["divider_ink"], bold=True, first=True, tracking=140)
            tf = self.box(s, 3.9, 3.7, 8.9, 2.8)
            self.para(tf, title, min(ctpt, 46), pal["ink"], bold=True, first=True, ls=1.05)
            if sub:
                self.para(tf, sub, 15, pal["sub"], sb=12, ls=1.25)
            return

        # solid / dark (default): full brand-color cover, light text
        s = self.dslide(pal["divider_bg"], "cover")
        self.rect(s, cx + 0.08, 0.7, 0.9, 0.12, self._bar_on(pal["divider_bg"]))
        tf = self.box(s, cx + 0.04, 0.95, 11.5, 0.5)
        self.para(tf, ey, 13, pal["divider_ink"], bold=True, first=True, tracking=180, align=al)
        tf = self.box(s, cx, 3.9, 12.1, 2.6)
        self.para(tf, title, ctpt, pal["divider_ink"], bold=True, first=True, ls=1.05, align=al)
        if sub:
            self.para(tf, sub, 16, pal["divider_ink"], sb=14, ls=1.25, align=al)

    def toc(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp.get("title", "Contents"),
                    sp.get("subtitle", ""), sp.get("page", ""))
        items = sp["items"]
        n = len(items); top = 3.0; bottom = 6.3
        step = (bottom - top) / max(n - 1, 1)
        for i, row in enumerate(items):
            num, title = row[0], row[1]
            pg = row[2] if len(row) > 2 else ""
            y = top + step * i
            self.rect(s, 0.62, y + 0.06, 11.9, 0.010, pal["border"])
            tf = self.box(s, 0.62, y, 0.9, 0.5)
            self.para(tf, num, 17, self.accent_ink, bold=True, first=True)
            tf = self.box(s, 1.7, y, 9.5, 0.5)
            self.para(tf, title, 16, pal["ink"], bold=True, first=True)
            if pg:
                tf = self.box(s, 11.7, y, 0.9, 0.5)
                self.para(tf, pg, 13, pal["muted"], first=True, align=PP_ALIGN.RIGHT)

    def divider(self, sp):
        pal, t = self.p, self.t
        style = t["divider_style"]
        num, title, sub = sp.get("num", ""), sp["title"], sp.get("subtitle")
        if sp.get("bg_image") and os.path.exists(sp["bg_image"]):
            self.imagefull({"image": sp["bg_image"], "eyebrow": num, "title": title, "subtitle": sub})
            return
        if style == "minimal":  # light bg, ink, thin section number — airy
            bg = sp.get("bg", pal["bg"])
            s = self.dslide(bg, "divider")
            self.rect(s, 0.62, 3.05, 0.9, 0.09, pal["accent"])
            tf = self.box(s, 0.62, 3.3, 12, 3.2)
            self.para(tf, num, 22, self.accent_ink, bold=True, first=True, tracking=60)
            self.para(tf, title, 36, pal["ink"], bold=True, sb=6)
            if sub:
                self.para(tf, sub, 15, pal["sub"], sb=10, ls=1.3)
            return
        bg = sp.get("bg", pal["divider_bg"])
        s = self.dslide(bg, "divider")
        ink = pal["divider_ink"]
        if style == "huge":  # oversized chapter numeral filling the slide
            tf = self.box(s, 0.5, 0.4, 12.4, 5.2)
            self.para(tf, num, 220, ink, bold=True, first=True, ls=0.9)
            tf = self.box(s, 0.62, 6.0, 12, 1.2)
            self.para(tf, title, 34, ink, bold=True, first=True)
            return
        self.rect(s, 0.7, 3.05, 0.9, 0.12, self._bar_on(bg))
        tf = self.box(s, 0.62, 3.3, 12, 3.2)
        self.para(tf, num, t["divider_num_pt"], ink, bold=True, first=True)
        self.para(tf, title, 38, ink, bold=True, sb=2)
        if sub:
            self.para(tf, sub, 15, ink, sb=10, ls=1.3)

    def _icon(self, s, name, x, y, size, color):
        """Draw a Lucide line icon; fall back to an accent chip if unavailable."""
        try:
            import icons
            png = icons.icon_png(name, color, int(size * 190), self.icons_dir)
        except Exception:
            png = None
        if png:
            s.shapes.add_picture(png, Inches(x), Inches(y), Inches(size), Inches(size))
            return True
        return False

    def icongrid(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        cells = sp["cells"]; cols = sp.get("cols", 2)
        rows = (len(cells) + cols - 1) // cols
        gx0, gy0 = 0.62, 2.92
        gw = (12.1 - (cols - 1) * 0.6) / cols
        gy_step = (6.75 - gy0) / rows
        for i, cell in enumerate(cells):
            has_icon = len(cell) >= 3
            icon = cell[0] if has_icon else None
            head, body = (cell[1], cell[2]) if has_icon else (cell[0], cell[1])
            r, c = divmod(i, cols)
            x = gx0 + c * (gw + 0.6)
            y = gy0 + r * gy_step
            drew = self._icon(s, icon, x, y, 0.44, pal["ink"]) if icon else False
            if not drew:
                self.rect(s, x, y, 0.46, 0.46, pal["accent"])
                tf = self.box(s, x, y, 0.46, 0.46, anchor=MSO_ANCHOR.MIDDLE)
                self.para(tf, str(i + 1), 14, self.on_accent, bold=True, first=True,
                          align=PP_ALIGN.CENTER)
            tf = self.box(s, x, y + 0.62, gw, 0.45)
            self.para(tf, head, 14.5, pal["ink"], bold=True, first=True)
            tf = self.box(s, x, y + 1.08, gw, gy_step - 1.25)
            self.para(tf, body, 12, pal["sub"], first=True, ls=1.4)

    def textfigure(self, sp):
        """Left column of head/body blocks + right zone reserved for an image
        (placed by place_images with slide=this, box in the right half)."""
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        items = sp["items"]
        y = 2.95
        step = min(1.35, (6.7 - y) / len(items))
        for head, body in items:
            tf = self.box(s, 0.62, y, 5.6, 0.42)
            self.para(tf, head, 15, pal["ink"], bold=True, first=True)
            tf = self.box(s, 0.62, y + 0.44, 5.6, step - 0.5)
            self.para(tf, body, 12, pal["sub"], first=True, ls=1.4)
            y += step
        if sp.get("image"):  # optional inline png (right half)
            self.pic_contain(s, sp["image"], 6.95, 3.0, 5.7, 3.75)
        elif sp.get("panel", True) and not sp.get("image_zone"):
            # native right panel when no image is placed
            px, py, pw, ph = 6.95, 2.98, 5.75, 3.7
            self.rect(s, px, py, pw, ph, pal["surface"])
            self.rect(s, px, py, 0.09, ph, pal["accent"])
            tf = self.box(s, px + 0.42, py + 0.36, pw - 0.8, 0.4)
            self.para(tf, sp.get("panel_title", "핵심 지향"), 12,
                      self.accent_ink, bold=True, first=True, tracking=60)
            tf = self.box(s, px + 0.42, py + 0.86, pw - 0.8, 1.4)
            self.para(tf, sp.get("panel_lede", sp.get("subtitle", "")), 15,
                      pal["ink"], bold=True, first=True, ls=1.3)
            pts = sp.get("panel_points", [])
            ty = py + ph - 0.28 - 0.42 * len(pts)
            for p in pts:
                tf = self.box(s, px + 0.42, ty, pw - 0.8, 0.4)
                self.para(tf, "· " + p, 12, pal["sub"], first=True)
                ty += 0.42

    def table(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        self._n += 1
        png = os.path.join(self.tmp, f"table_{self._n}.png")
        try:
            import charts
            charts.comparison_table(sp["headers"], sp["rows"], png, ink=pal["ink"],
                                    sub=pal["sub"], muted=pal["muted"], accent=pal["accent"],
                                    line=pal["border"], panel=pal["surface"])
            self.pic_contain(s, png, 0.62, 2.7, 12.1, 4.0)
        except Exception as e:
            sys.stderr.write(f"[table degraded: {e}]\n")

    def numbered(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        items = sp["items"]
        lede = sp.get("lede")
        x = 4.4 if lede else 0.62
        if lede:
            tf = self.box(s, 0.62, 2.95, 3.3, 3.6)
            self.para(tf, lede, 12.5, pal["sub"], first=True, ls=1.5)
        y = 2.9
        step = (6.75 - y) / len(items)
        for i, it in enumerate(items):
            head, body = it[0], it[1]
            yy = y + step * i
            tf = self.box(s, x, yy, 0.95, 0.6)
            self.para(tf, f"{i + 1:02d}", 21, self.accent_ink, bold=True, first=True)
            tf = self.box(s, x + 1.0, yy, 12.7 - x - 1.0, 0.42)
            self.para(tf, head, 14, pal["ink"], bold=True, first=True)
            tf = self.box(s, x + 1.0, yy + 0.42, 12.7 - x - 1.0, step - 0.5)
            self.para(tf, body, 11.5, pal["sub"], first=True, ls=1.32)

    def bullets(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        items = sp["items"]; x = 0.62; y = 2.85
        step = (6.8 - y) / len(items)
        for i, it in enumerate(items):
            head, body = it[0], it[1]
            yy = y + step * i
            tf = self.box(s, x, yy, 0.9, 0.6)
            self.para(tf, f"{i + 1:02d}", 22, self.accent_ink, bold=True, first=True)
            tf = self.box(s, x + 0.95, yy, 11.4, 0.45)
            self.para(tf, head, 15, pal["ink"], bold=True, first=True)
            tf = self.box(s, x + 0.95, yy + 0.46, 11.4, step - 0.5)
            self.para(tf, body, 11.5, pal["sub"], first=True, ls=1.32)

    def kpi(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        self._n += 1
        png = os.path.join(self.tmp, f"kpi_{self._n}.png")
        drawn = False
        try:
            import charts
            charts.kpi_bars(sp["metrics"], png, ink=pal["ink"], sub=pal["sub"],
                            muted=pal["muted"], accent=pal["accent"], line=pal["border"])
            drawn = os.path.exists(png)
        except Exception as e:
            sys.stderr.write(f"[kpi chart degraded to text: {e}]\n")
        if drawn:
            self.pic_contain(s, png, 0.62, 2.75, 12.1, 3.9)
        else:
            items = [(m[0], f"{m[1]} → {m[2]}") for m in sp["metrics"]]
            self.bullets({"title": sp["title"], "eyebrow": sp.get("eyebrow", ""),
                          "subtitle": sp.get("subtitle", ""), "items": items,
                          "page": sp.get("page", "")})
            return
        if sp.get("note"):
            tf = self.box(s, 0.62, 6.75, 12.1, 0.4)
            self.para(tf, sp["note"], 10.5, pal["muted"], first=True)

    def roadmap(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        png = os.path.join(self.tmp, "roadmap.png")
        try:
            import charts
            phases = [tuple(p) for p in sp["phases"]]
            charts.roadmap(phases, png, ink=pal["ink"], sub=pal["sub"], muted=pal["muted"],
                           accent=pal["accent"], line=pal["border"], panel=pal["surface"])
            self.pic_contain(s, png, 0.62, 2.7, 12.1, 4.0)
        except Exception as e:
            sys.stderr.write(f"[roadmap degraded: {e}]\n")

    def bignum(self, sp):
        """A row of big-number KPI cards (value + label + delta). For financial highlights."""
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        cards = sp["cards"]
        n = len(cards)
        gx0 = 0.62
        gw = (12.1 - (n - 1) * 0.4) / n
        y, h = 3.05, 2.7
        for i, c in enumerate(cards):
            val, label = c[0], c[1]
            delta = c[2] if len(c) > 2 else None
            x = gx0 + i * (gw + 0.4)
            self.rect(s, x, y, gw, h, pal["surface"])
            self.rect(s, x, y, gw, 0.08, pal["accent"])
            if self.t["card_icons"]:
                ic = resolve_icon(label)
                if ic:
                    self._icon(s, ic, x + gw - 0.64, y + 0.3, 0.36, pal["accent"])
            tf = self.box(s, x + 0.3, y + 0.4, gw - 0.9, 0.4)
            self.para(tf, label, 12, pal["sub"], first=True)
            tf = self.box(s, x + 0.3, y + 0.92, gw - 0.6, 1.0)
            self.para(tf, val, 33, pal["ink"], bold=True, first=True)
            if delta:
                tf = self.box(s, x + 0.3, y + h - 0.62, gw - 0.6, 0.4)
                self.para(tf, delta, 12.5, self.accent_ink, bold=True, first=True)

    def trend(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        self._n += 1
        png = os.path.join(self.tmp, f"trend_{self._n}.png")
        try:
            import charts
            charts.trend_dual(sp["labels"], sp["revenue"], sp["op"], png, ink=pal["ink"],
                              sub=pal["sub"], muted=pal["muted"], accent=pal["accent"],
                              line=pal["border"], rev_name=sp.get("rev_name", "매출"),
                              op_name=sp.get("op_name", "영업이익"), unit=sp.get("unit", "조원"))
            self.pic_contain(s, png, 0.62, 2.7, 12.1, 4.0)
        except Exception as e:
            sys.stderr.write(f"[trend degraded: {e}]\n")

    def segment(self, sp):
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""),
                    sp.get("page", ""))
        self._n += 1
        png = os.path.join(self.tmp, f"seg_{self._n}.png")
        try:
            import charts
            charts.hbars(sp["labels"], sp["values"], png, ink=pal["ink"], sub=pal["sub"],
                         muted=pal["muted"], accent=pal["accent"], line=pal["border"],
                         unit=sp.get("unit", "조원"))
            self.pic_contain(s, png, 0.62, 2.75, 12.1, 3.9)
        except Exception as e:
            sys.stderr.write(f"[segment degraded: {e}]\n")

    def imagefull(self, sp):
        """Full-bleed image with a solid caption band + title overlay (image-forward)."""
        pal = self.p
        s = self.slide(pal["divider_bg"])
        self._img_or_placeholder(s, sp, 0, 0, SW, SH)
        self.rect(s, 0, SH - 2.4, SW, 2.4, pal["divider_bg"])
        self.rect(s, 0.7, SH - 2.05, 0.9, 0.09, self._bar_on(pal["divider_bg"]))
        if sp.get("eyebrow"):
            tf = self.box(s, 0.7, SH - 1.85, 11, 0.35)
            self.para(tf, sp["eyebrow"].upper(), 12, pal["divider_ink"], bold=True,
                      first=True, tracking=160)
        tf = self.box(s, 0.66, SH - 1.5, 12, 1.2)
        self.para(tf, sp["title"], 38, pal["divider_ink"], bold=True, first=True, ls=1.05)
        if sp.get("subtitle"):
            self.para(tf, sp["subtitle"], 15, pal["divider_ink"], sb=8)

    def imagesplit(self, sp):
        """Half image (cover-fit) + half text (title + items)."""
        pal = self.p
        s = self.cslide()
        half = SW / 2
        side = sp.get("side", "right")
        ix = half if side == "right" else 0
        tx = 0.72 if side == "right" else half + 0.6
        self._img_or_placeholder(s, sp, ix, 0, half, SH)
        if sp.get("eyebrow"):
            tf = self.box(s, tx, 0.95, half - 1.3, 0.4)
            self.para(tf, sp["eyebrow"].upper(), 11, self.accent_ink, bold=True, first=True,
                      tracking=120)
        tf = self.box(s, tx, 1.45, half - 1.3, 1.4)
        self.para(tf, sp["title"], 26, pal["ink"], bold=True, first=True)
        self.rect(s, tx, 2.72, 0.9, 0.05, pal["accent"])
        y = 3.15
        for it in sp.get("items", []):
            if isinstance(it, (list, tuple)):
                tf = self.box(s, tx, y, half - 1.3, 0.42)
                self.para(tf, it[0], 15, pal["ink"], bold=True, first=True)
                tf = self.box(s, tx, y + 0.44, half - 1.3, 0.8)
                self.para(tf, it[1], 12, pal["sub"], first=True, ls=1.35)
                y += 1.3
            else:
                tf = self.box(s, tx, y, half - 1.3, 0.5)
                self.para(tf, "· " + it, 13, pal["sub"], first=True)
                y += 0.56

    def imagegrid(self, sp):
        """A grid of images with captions."""
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""), sp.get("page", ""))
        imgs = sp["images"]
        cols = sp.get("cols", 2)
        rows = (len(imgs) + cols - 1) // cols
        gx0, gy0 = 0.62, 2.85
        gw = (12.1 - (cols - 1) * 0.4) / cols
        gh = (6.9 - gy0) / rows - 0.35
        for i, im in enumerate(imgs):
            r, c = divmod(i, cols)
            x = gx0 + c * (gw + 0.4)
            y = gy0 + r * (gh + 0.55)
            self._img_or_placeholder(s, {"image": im.get("image"), "title": im.get("caption", "")},
                                     x, y, gw, gh)
            if im.get("caption"):
                tf = self.box(s, x, y + gh + 0.06, gw, 0.4)
                self.para(tf, im["caption"], 12, pal["sub"], first=True)

    def diagram(self, sp):
        """Architecture / process diagram. `flow`: auto horizontal boxes + arrows.
        `nodes`+`edges`: explicit boxes at [label,x,y,w,h] + connectors [i,j]."""
        pal = self.p
        s = self.cslide()
        self.header(s, sp.get("eyebrow", ""), sp["title"], sp.get("subtitle", ""), sp.get("page", ""))
        if sp.get("flow"):
            boxes = sp["flow"]
            n = len(boxes)
            gap = 0.5
            bw = (12.1 - gap * (n - 1)) / n
            bh, y = 2.0, 3.7
            for i, lab in enumerate(boxes):
                x = 0.62 + i * (bw + gap)
                self.rect(s, x, y, bw, bh, pal["surface"])
                self.rect(s, x, y, bw, 0.09, pal["accent"])
                tf = self.box(s, x + 0.2, y, bw - 0.4, bh, anchor=MSO_ANCHOR.MIDDLE)
                if isinstance(lab, (list, tuple)):
                    self.para(tf, lab[0], 15, pal["ink"], bold=True, first=True, align=PP_ALIGN.CENTER)
                    self.para(tf, lab[1], 11.5, pal["sub"], sb=5, align=PP_ALIGN.CENTER, ls=1.3)
                else:
                    self.para(tf, lab, 15, pal["ink"], bold=True, first=True, align=PP_ALIGN.CENTER)
                if i < n - 1:
                    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.06),
                                            Inches(y + bh / 2 - 0.14), Inches(gap - 0.12), Inches(0.28))
                    ar.fill.solid(); ar.fill.fore_color.rgb = C(pal["accent"])
                    ar.line.fill.background(); ar.shadow.inherit = False
        elif sp.get("nodes"):
            centers = []
            for nd in sp["nodes"]:
                lab, x, y, w, h = nd
                self.rect(s, x, y, w, h, pal["surface"])
                self.rect(s, x, y, 0.08, h, pal["accent"])
                tf = self.box(s, x + 0.22, y, w - 0.32, h, anchor=MSO_ANCHOR.MIDDLE)
                self.para(tf, lab, 13, pal["ink"], bold=True, first=True)
                centers.append((x + w / 2, y + h / 2))
            for e in sp.get("edges", []):
                (x1, y1), (x2, y2) = centers[e[0]], centers[e[1]]
                cn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
                cn.line.color.rgb = C(pal["accent"]); cn.line.width = Pt(1.5)

    def closing(self, sp):
        pal = self.p
        s = self.dslide(pal["divider_bg"], "cover")
        if self.t["closing_mark"]:
            ic = resolve_icon(sp.get("title", "")) or "flag"
            self._icon(s, ic, 0.62, 2.0, 0.62, self._bar_on(pal["divider_bg"]))
        self.rect(s, 0.7, 2.9, 0.9, 0.12, self._bar_on(pal["divider_bg"]))
        tf = self.box(s, 0.62, 3.15, 12.1, 2.6)
        self.para(tf, sp.get("title", "Thank you"), 40, pal["divider_ink"], bold=True, first=True)
        if sp.get("subtitle"):
            self.para(tf, sp["subtitle"], 15, pal["divider_ink"], sb=12, ls=1.35)
        if sp.get("contact"):
            tf = self.box(s, 0.62, 6.6, 12.1, 0.4)
            self.para(tf, sp["contact"], 11, pal["divider_ink"], first=True)

    def build(self, spec):
        dispatch = {"cover": self.cover, "toc": self.toc, "divider": self.divider,
                    "icongrid": self.icongrid, "kpi": self.kpi, "bullets": self.bullets,
                    "roadmap": self.roadmap, "closing": self.closing,
                    "textfigure": self.textfigure, "table": self.table,
                    "numbered": self.numbered, "bignum": self.bignum,
                    "trend": self.trend, "segment": self.segment,
                    "imagefull": self.imagefull, "imagesplit": self.imagesplit,
                    "imagegrid": self.imagegrid, "diagram": self.diagram}
        for sl in spec["slides"]:
            fn = dispatch.get(sl["layout"])
            if not fn:
                raise ValueError(f"unknown layout: {sl['layout']}")
            fn(sl)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--palette", required=True)
    ap.add_argument("--spec", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--pdf", action="store_true", help="also export PDF via soffice")
    ap.add_argument("--theme", help="visual theme name (from _engine/themes.json) or a .json path")
    args = ap.parse_args()

    with open(args.palette, encoding="utf-8") as f:
        pal = json.load(f)
    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)

    # resolve theme: --theme flag > spec.meta.theme; a name looks up themes.json
    theme = None
    theme_ref = args.theme or spec.get("meta", {}).get("theme")
    if isinstance(theme_ref, dict):
        theme = theme_ref
    elif theme_ref:
        if os.path.isfile(theme_ref):
            with open(theme_ref, encoding="utf-8") as f:
                theme = json.load(f)
        else:
            tf = os.path.join(os.path.dirname(os.path.abspath(__file__)), "themes.json")
            if os.path.isfile(tf):
                with open(tf, encoding="utf-8") as f:
                    theme = json.load(f).get("themes", {}).get(theme_ref)
                if theme is None:
                    sys.stderr.write(f"[theme '{theme_ref}' not found; using default]\n")

    # charts.py sits next to this file
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

    os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        deck = Deck(pal, tmp, theme)
        deck.build(spec)
        deck.prs.save(args.out)
    print(f"[ok] {args.out}  ({len(spec['slides'])} slides, palette={pal.get('name')})")

    if args.pdf:
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            sys.stderr.write("[pdf skipped: soffice not found]\n")
            return
        outdir = os.path.dirname(os.path.abspath(args.out))
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir",
                        outdir, args.out], check=True,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        pdf = os.path.splitext(args.out)[0] + ".pdf"
        print(f"[ok] {pdf}")


if __name__ == "__main__":
    main()
