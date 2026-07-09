#!/usr/bin/env python3
"""place_images.py -- OPTIONAL hero-image generator for a rendered deck.

Adds AI-generated images to slides where you ask for them. Two paths, auto-selected:
  Path A (OPENAI_API_KEY set + openai lib): generate each image and fit it in the box.
  Path B (no key / failure): draw an X-box placeholder carrying the prompt text.
So this never blocks a build -- without a key you still get a laid-out placeholder.

Image plan JSON:
{
  "style_suffix": ", clean minimal vector illustration, brand-colored, white background, no text, no logo, no watermark",
  "size": "1024x1024", "quality": "high",
  "images": [
    {"slide": 1, "box_in": [7.7, 3.6, 5.2, 3.4], "prompt": "secure AI datacenter with memory chips"}
  ]
}
box_in = [x, y, w, h] in inches (slide is 13.333 x 7.5). slide is 1-based.

Usage:
  python3 place_images.py --in deck.pptx --plan images.json --out deck_img.pptx \
      [--palette palette.json] [--mode auto|generate|placeholder] [--model gpt-image-1]

Setup for real image generation:
  export OPENAI_API_KEY=sk-...        # your own key; the only outbound call in this repo
  pip install openai pillow
  (model configurable via --model or IMAGE_MODEL env; default gpt-image-1)
"""
import argparse
import base64
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _c(h):
    return RGBColor.from_string(h.lstrip("#").upper())


def gen_openai(prompt, size, quality, model, out_png):
    """Return True if an image was generated to out_png, else False."""
    key = os.environ.get("OPENAI_API_KEY")
    if not key:
        return False
    try:
        from openai import OpenAI
    except Exception:
        sys.stderr.write("[images] openai lib missing -> placeholder\n")
        return False
    try:
        client = OpenAI(api_key=key)
        r = client.images.generate(model=model, prompt=prompt, size=size,
                                    quality=quality, n=1)
        b64 = r.data[0].b64_json
        with open(out_png, "wb") as f:
            f.write(base64.b64decode(b64))
        return True
    except Exception as e:
        sys.stderr.write(f"[images] generate failed ({e}) -> placeholder\n")
        return False


def contain(png, x, y, w, h):
    from PIL import Image
    iw, ih = Image.open(png).size
    bx, ix = w / h, iw / ih
    if ix > bx:
        nw, nh = w, w / ix
    else:
        nh, nw = h, h * ix
    return x + (w - nw) / 2, y + (h - nh) / 2, nw, nh


def placeholder(slide, x, y, w, h, prompt, ink, surface, border):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = _c(surface)
    box.line.color.rgb = _c(border); box.line.width = Pt(1)
    box.shadow.inherit = False
    ln = slide.shapes.add_connector  # draw an X with two lines
    for a, b in (((x, y), (x + w, y + h)), ((x + w, y), (x, y + h))):
        c = ln(2, Inches(a[0]), Inches(a[1]), Inches(b[0]), Inches(b[1]))
        c.line.color.rgb = _c(border); c.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"[image]\n{prompt}"
    r.font.size = Pt(9); r.font.color.rgb = _c(ink)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="inp", required=True)
    ap.add_argument("--plan", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--palette")
    ap.add_argument("--mode", default="auto", choices=["auto", "generate", "placeholder"])
    ap.add_argument("--model", default=os.environ.get("IMAGE_MODEL", "gpt-image-1"))
    args = ap.parse_args()

    with open(args.plan, encoding="utf-8") as f:
        plan = json.load(f)
    pal = {"ink": "1A1A1A", "surface": "F5F6F8", "border": "E1E3E8"}
    if args.palette:
        with open(args.palette, encoding="utf-8") as f:
            pal.update({k: v for k, v in json.load(f).items() if k in pal})

    prs = Presentation(args.inp)
    slides = list(prs.slides)
    suffix = plan.get("style_suffix", "")
    size = plan.get("size", "1024x1024")
    quality = plan.get("quality", "high")
    tmpdir = os.path.join(os.path.dirname(os.path.abspath(args.out)) or ".", "assets")
    os.makedirs(tmpdir, exist_ok=True)

    gen = ph = 0
    for i, im in enumerate(plan.get("images", [])):
        s = slides[im["slide"] - 1]
        x, y, w, h = im["box_in"]
        prompt = im["prompt"] + im.get("style_suffix", suffix)
        want_gen = args.mode in ("auto", "generate")
        png = os.path.join(tmpdir, f"img_{i}.png")
        done = False
        if want_gen:
            done = gen_openai(prompt, im.get("size", size), im.get("quality", quality),
                              args.model, png)
        if done:
            cx, cy, cw, ch = contain(png, x, y, w, h)
            s.shapes.add_picture(png, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
            gen += 1
        else:
            if args.mode == "generate":
                sys.stderr.write("[images] --mode generate but no key/failure; placeholder\n")
            placeholder(s, x, y, w, h, im["prompt"], pal["ink"], pal["surface"], pal["border"])
            ph += 1

    prs.save(args.out)
    print(json.dumps({"out": args.out, "generated": gen, "placeholders": ph,
                      "mode": args.mode, "model": args.model}, ensure_ascii=False))


if __name__ == "__main__":
    main()
